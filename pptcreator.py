import re
from pptx import Presentation   
from slidenum import slidenum,get_indices,get_titleindices
from dataextractor import substitute_text,extractdata
from notes import snaddition,snextraction
import os 

def pptcreate(excel,ppt,out,sninput):
    speaker_notes=extractdata(inputfile=sninput)
    # print(speaker_notes)
    # for k, v in speaker_notes.items():
    #     if 'header' not in v or 'text' not in v:
    #         print(f"Missing key in entry {k}: {v}")
    head=[]
    text=[]
    file_path=sninput
    fname=os.path.splitext(os.path.basename(file_path))[0]
    split=fname.split('_')
    title=split[-1]
    prs=Presentation(ppt)
    slide=prs.slide_layouts[1]
    addslide=prs.slides.add_slide(slide)
    for shape in addslide.placeholders:
        if shape.name=='Title 1':
            shape.text=title 
            
            
    for i in speaker_notes.keys():
        if i == 1:
            continue
        if speaker_notes[i]['text']==speaker_notes[i]['header']:
            speaker_notes[i]['text']='Image'
        ost = substitute_text(speaker_notes[i]['text'], speaker_notes[i]['header'], "")
        ost = re.sub(r'\n+', '\n', ost)
        speaker_notes[i]['text']=ost
        text.append(speaker_notes[i]['text'])
        head.append(speaker_notes[i]['header'].strip())
    # print(speaker_notes)        

    c=0
    slidenumber=0
    prev=[]
    for i in text:
        header=head[text.index(i)]
        # print(header)
        # print(header)
        if header.lower() not in ['learning objectives','summary','overview','introduction']:
            if i!='Image':
                print(f"\n{header}")
                c+=1
                if i=='':
                    continue
                split=i.split("\n")
                for i in split:
                    if i=='':
                        split.remove(i)
                bps=len(split)
                if bps >6:          #Change according to max value in layout_data.xlsx  
                    bps=1
                    split=['\n•'.join(split)]
                elif bps==0:
                    bps=1
                    split=['NA']
                    
                    
                # print(split,'\n')
                print('Number of Bullet points :',bps)
                slider=slidenum(bps,excel,ppt)
                if slider not in prev:
                    prev.append(slider)
                else:
                    slider=slidenum(bps,excel,ppt)
                    prev.append(slider)
  
                idxl=get_indices(excel,slider)
                titl=get_titleindices(excel,slider)
                slidelay=prs.slide_layouts[slider]
                slide=prs.slides.add_slide(slidelay)
                slidenumber+=1
                
                print(f'PPT Slide number :{slidenumber}')
                print(f'Slide layout number:{slider}')
                # print(f"idxs : {idxl}")       Uncomment to check the idx vals
                
                
                for shape in slide.placeholders:
                    
                    if shape.placeholder_format.idx in titl:
                        print("✅Identified Header")
                        shape.text=header
                        print('✅Added Header')
                        prs.save(out)
                        
                    if shape.placeholder_format.idx in idxl and split:
                        print("✅Identified Textbox")
                        shape.text=split[0]
                        idxl.pop(0)
                        split.pop(0)
                        print('✅Added Text')
                        prs.save(out)
                    
            else:
                for slide in prs.slide_layouts:
                    if slide.name=='H + SH + (F) PIC':
                        addslide=prs.slides.add_slide(slide) 
                        for shape in addslide.placeholders:
                            if shape.name=='Heading':
                                shape.text=header

        else:
            if header.lower()=='learning objectives':
                print(f"✅Adding {header} slide")
                for slide in prs.slide_layouts:
                    if slide.name=='003_LO_ANI':
                        addslide=prs.slides.add_slide(slide)
                        for shape in addslide.shapes:
                            if shape.placeholder_format.idx==17:
                                print("✅Identified Textbox")
                                removal='By the end of this session, you will be able to:'
                                if removal in i:
                                    i=i.replace(removal,"")
                                    print("True")
                                print(i)
                                shape.text=i
                                print('✅Added Text')
                                prs.save(out)
            elif header.lower()=='summary':
                print(f"✅Adding {header} slide")
                for slide in prs.slide_layouts:
                    if slide.name=='Summary_With ANI':
                        addslide=prs.slides.add_slide(slide)
                        for shape in addslide.shapes:
                            if shape.placeholder_format.idx==17:
                                print("✅Identified Textbox")
                                removal='In this topic you have learnt:'
                                if removal in i:
                                    i=i.replace(removal,"")
                                shape.text=i    
                                print('✅Added Text')              
                                prs.save(out)
            elif header.lower()=='overview':
                print(f"✅Adding {header} slide")
                for slide in prs.slide_layouts:
                    if slide.name=='H + (F) VID + ML':
                        addslide=prs.slides.add_slide(slide)
                        for shape in addslide.shapes:
                            if shape.placeholder_format.idx==10:
                                print("✅Identified Header")
                                shape.text=header
                                print('✅Added Header')
                                prs.save(out)
                            if shape.placeholder_format.idx==13:
                                print("✅Identified Textbox")
                                shape.text=i      
                                print('✅Added Text')                
                                prs.save(out)
            elif header.lower()=='introduction':
                print(f"✅Adding {header} slide")
                for slide in prs.slide_layouts:
                    if slide.name=='H + (R) PIC + 3/4 TXT':
                        addslide=prs.slides.add_slide(slide)
                        print(slide.name)
                        for shape in addslide.shapes:
                            if shape.placeholder_format.idx==10:
                                print("✅Identified Header")
                                shape.text=header
                                print('✅Added Header')
                                prs.save(out)
                            if shape.placeholder_format.idx==15:
                                print("✅Identified Textbox")
                                shape.text=i      
                                print('✅Added Text')                
                                prs.save(out)



    try:
        notel=snextraction(sninput)
        counter=0
        for i in notel:
            snaddition(out,i,counter)
            counter+=1
        print("✅Notes for slides added")
    
    except:  # noqa: E722
        print("Unable to add all speaker notes ")
# pptcreate(excel,ppt,out,sninput)