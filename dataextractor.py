from pptx import Presentation
import re



def substitute_text(text, pattern, replacement):
    """
    Substitutes text from a string with " ".

    Args:
        text: The input string.
        pattern: The regular expression pattern to match.
        replacement: The replacement string.

    Returns:
        The modified string with the matched text replaced.
    """
    if text !='':
      return re.sub(re.escape(pattern), replacement, text)
    else:
      text='NOT AVAILABLE'
      return text

def remove_extra_whitespaces(text):
  """Removes leading, trailing, and extra whitespaces within a string.

  Args:
    text: The input string.

  Returns:
    The string with extra whitespaces removed.
  """
  text = text.strip()  # Remove leading/trailing whitespaces
  text = re.sub(r'\s+', ' ', text)  # Replace multiple whitespaces with single space
  return text

def extract_top_textbox_text(slide):
  """
  Checks if a slide in a presentation has a textbox at the top and extracts its text.

  Args:
    prs_path: Path to the presentation file.

  Returns:
    The text content of the topmost textbox on the first slide, or None if no such textbox exists.
  """
  try:

    # Find the topmost textbox
    topmost_textbox = None
    min_top = float('inf')

    for shape in slide.shapes:
      if shape.shape_type == 17 and shape.text:  # Check for textboxes with text
        top = shape.top
        if top < min_top:
          min_top = top
          topmost_textbox = shape

    if topmost_textbox:
      return topmost_textbox.text
    else:
      return None
  except Exception as e:
      print(f"An error occurred: {e}")
      return None

def extractdata(inputfile):
  try :
      prs = Presentation(inputfile)
      speaker_notes = {}
      for j, slide in enumerate(prs.slides):

        #print(j)
        speaker_notes[j+1] = {}

        text_on_slide = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text :
                  text_on_slide += text_frame.text
        speaker_notes[j+1]['text'] = text_on_slide

        notes_slide = slide.notes_slide
        if notes_slide:
          text_frame = notes_slide.notes_text_frame
          if text_frame.text:
            speaker_notes[j+1]['speaker_notes'] = text_frame.text

        # Find the header (assuming first line of text is the header)
        if extract_top_textbox_text(slide):
          speaker_notes[j+1]['header'] = extract_top_textbox_text(slide)
          header_found = True
        else :
          speaker_notes[j+1]['header']='Header absent in source'
          header_found = False
          for shape in slide.shapes:
            if shape.has_text_frame:
              text_frame = shape.text_frame
              if text_frame.text:
                speaker_notes[j+1]['header'] = text_frame.text.split("\n")[0]
                header_found = True  # noqa: F841
                #speaker_notes[j+1]['text'] = "\n".join(text_frame.text.split("\n")[1:]).strip()
                break


    
  except Exception as e:
      print(f"An error occurred: {e}")
      print("An error occured at : Slide")

  return speaker_notes
