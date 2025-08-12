# prompt: From the excel sheet count the number of 'Text Placeholder 2' in shape name per layout and return a dictionary

import pandas as pd
import random
from pptx import Presentation
def count_placeholders_per_layout(excel_file):
    """
    Counts the number of shapes with the name 'Text Placeholder 2' per layout
    in an Excel sheet.

    Args:
        excel_file (str): The path to the Excel file.

    Returns:
        dict: A dictionary where keys are layout names and values are
              the counts of 'Text Placeholder 2' shapes in that layout.
    """
    # Assuming the Excel file has a column named 'Layout' and a column named 'Shape Name'
    # You might need to adjust the sheet_name and column names based on your Excel file structure.
    try:
        df = pd.read_excel(excel_file)
    except FileNotFoundError:
        return {"error": f"File not found at {excel_file}"}

    if 'Layout number' not in df.columns or 'Shape Name' not in df.columns:
        return {"error": "Excel file must contain 'Layout' and 'Shape Name' columns."}

    # Filter rows where 'Shape Name' is 'Text Placeholder 2'
    placeholder_df = df[df['Shape Name'] == 'Text Placeholder 2']

    # Group by 'Layout' and count the occurrences
    layout_counts = placeholder_df['Layout number'].value_counts().to_dict()

    final=dict(sorted(layout_counts.items()))
    return final


# Example usage:
# Assuming you have an Excel file named 'your_excel_file.xlsx'

def slidenum(bps,excel,layppt):
    prs=Presentation(layppt)
    final=count_placeholders_per_layout(excel)
    # print(final)
    l=[]  # noqa: E741
    for i in final.keys():
        if final[i]==bps:
          slidelay=prs.slide_layouts[i]
          if slidelay.name.lower() not in ['learning objectives','summary','overview','introduction','image','summary_without_animation','instructor intro_1_animated']:  
            l.append(i)
    random.shuffle(l)            
    choice=random.choice(l)
    return choice
    

def get_indices(excel_file, chosen_layout):
  """
  Retrieves the shape indices of all Text Placeholder 2 shapes for a chosen layout.

  Args:
      excel_file (str): The path to the Excel file.
      chosen_layout (str or int): The name or number of the chosen layout.

  Returns:
      list: A list of shape indices for the Text Placeholder 2 shapes
            in the chosen layout, or an empty list if the layout is not found
            or no such shapes exist.
  """
  try:
    df = pd.read_excel(excel_file)
  except FileNotFoundError:
    print(f"Error: File not found at {excel_file}")
    return []

  if 'Layout number' not in df.columns or 'Shape Name' not in df.columns or 'Shape idx' not in df.columns:
    print("Error: Excel file must contain 'Layout number ', 'Shape Name', and 'Shape Idx' columns.")
    return []

  # Filter by the chosen layout
  layout_df = df[df['Layout number'] == chosen_layout]

  # Filter by 'Shape Name' which is 'Text Placeholder 2'
  placeholder2_df = layout_df[layout_df['Shape Name'] == 'Text Placeholder 2']

  # Extract the 'Shape Index' and convert to a list
  shape_indices = placeholder2_df['Shape idx'].tolist()

  return shape_indices



def get_titleindices(excel_file, chosen_layout):
  """
  Retrieves the shape indices of all Text Placeholder 2 shapes for a chosen layout.

  Args:
      excel_file (str): The path to the Excel file.
      chosen_layout (str or int): The name or number of the chosen layout.

  Returns:
      list: A list of shape indices for the Text Placeholder 2 shapes
            in the chosen layout, or an empty list if the layout is not found
            or no such shapes exist.
  """
  try:
    df = pd.read_excel(excel_file)
  except FileNotFoundError:
    print(f"Error: File not found at {excel_file}")
    return []

  if 'Layout number' not in df.columns or 'Shape Name' not in df.columns or 'Shape idx' not in df.columns:
    print("Error: Excel file must contain 'Layout number ', 'Shape Name', and 'Shape Idx' columns.")
    return []

  # Filter by the chosen layout
  layout_df = df[df['Layout number'] == chosen_layout]

  # Filter by 'Shape Name' which is 'Text Placeholder 2'
  placeholder2_df = layout_df[layout_df['Shape Name'] == 'Heading']

  # Extract the 'Shape Index' and convert to a list
  shape_indices = placeholder2_df['Shape idx'].tolist()

  return shape_indices


# print(count_placeholders_per_layout(r'C:\Users\Siddharth Bhat\Downloads\layout_shapes.xlsx'))

# a=slidenum(1)
# print(a)