from pptx import Presentation
def snextraction(path):
    l=[]  # noqa: E741
    # Load your presentation
    prs = Presentation(path)

    # List to store notes for each slide
    notes = []

    for idx, slide in enumerate(prs.slides):
        # Check if the slide has a notes slide
        if slide.has_notes_slide:
            note_text = slide.notes_slide.notes_text_frame.text
        else:
            note_text = ""
        notes.append((idx, note_text))

    # Print notes for each slide
    for slide_num, note in notes:
        # print(f"Slide {slide_num + 1} notes: {note}")
        l.append(note)
    return l

def snaddition(path,text,idx):
    # Load your presentation
    prs = Presentation(path)

    # Select the slide you want to add notes to (0-based index)
    slide = prs.slides[idx]

    # Access the notes slide; if it doesn't exist, python-pptx will create it
    notes_slide = slide.notes_slide

    # Set the notes text
    notes_slide.notes_text_frame.text = text

    # Save the presentation
    prs.save(path)
