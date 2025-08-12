from pptx import Presentation
import pandas as pd

def analyze_ppt(ppt_path, excel_path):
    """
    Analyzes the number of text placeholders, headings, subheadings, and picture placeholders
    per layout in a PowerPoint presentation and saves the results to an Excel file.

    Args:
        ppt_path (str): The path to the PowerPoint file.
        excel_path (str): The path to save the Excel file.
    """
    prs = Presentation(ppt_path)
    slide_data = []
    counter=0
    for slide in prs.slide_layouts:
        counter+=1
        slide_info = {
            'Slide Number':counter ,
            'Layout Name': slide.name,
            'Text Placeholders': 0,
            'Heading Placeholders': 0,
            'Subheading Placeholders': 0,
            'Picture Placeholders': 0
        }

        for shape in slide.placeholders:
            # Check for text placeholders
            if shape.name.lower()=='heading':
                slide_info['Heading Placeholders']+=1
            elif shape.name.lower()=='subheading':
                slide_info['Subheading Placeholders']+=1
            elif shape.name.lower()=='text placeholder 2':
                slide_info['Text Placeholders']+=1
            elif 'picture placeholder' in shape.name.lower():
                slide_info['Picture Placeholders']+=1

        slide_data.append(slide_info)

    df = pd.DataFrame(slide_data)

    # You can also aggregate data by layout if needed
    # layout_summary = df.groupby('Layout Name').sum()

    df.to_excel(excel_path, index=False)
    print(f"Analysis complete. Results saved to {excel_path}")
    
    
    
analyze_ppt("layout_template.pptx", 'layout_data.xlsx')